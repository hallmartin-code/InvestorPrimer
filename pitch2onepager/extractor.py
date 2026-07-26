"""Deck ingestion: parse a .pdf or .pptx into structured text."""

from __future__ import annotations

from pathlib import Path

from .models import DeckContent, SlideContent
from .utils import (
    ExtractionError,
    FileError,
    clean_text,
    detect_file_type,
    guess_title,
    validate_input_file,
)


def extract_deck(file_path: str) -> DeckContent:
    """Parse a pitch deck into a :class:`DeckContent`.

    Args:
        file_path: Path to a ``.pdf`` or ``.pptx`` file.

    Raises:
        FileError: file missing or unsupported extension.
        ExtractionError: the file exists but could not be parsed at all.
    """
    path = validate_input_file(file_path)
    file_type = detect_file_type(path)

    if file_type == "pdf":
        slides = _extract_pdf(path)
    else:
        slides = _extract_pptx(path)

    if not slides:
        raise ExtractionError(
            f"No pages or slides found in {path.name}. The file may be empty or corrupt."
        )

    full_text = "\n\n".join(
        f"--- Slide {s.slide_number} ---\n{s.raw_text}" for s in slides if s.raw_text
    )

    return DeckContent(
        source_file=str(path),
        file_type=file_type,  # type: ignore[arg-type]
        slide_count=len(slides),
        slides=slides,
        full_text=full_text,
    )


# --------------------------------------------------------------------------- #
# PDF
# --------------------------------------------------------------------------- #


def _extract_pdf(path: Path) -> list[SlideContent]:
    from pypdf import PdfReader

    try:
        reader = PdfReader(str(path))
    except Exception as exc:  # noqa: BLE001 - pypdf raises a wide variety of errors
        raise ExtractionError(f"Could not open PDF '{path.name}': {exc}") from exc

    slides: list[SlideContent] = []
    for index, page in enumerate(reader.pages, start=1):
        try:
            raw = page.extract_text() or ""
        except Exception:  # noqa: BLE001 - a single malformed page must not abort the run
            raw = ""
        text = clean_text(raw)
        slides.append(
            SlideContent(
                slide_number=index,
                title=guess_title(text),
                raw_text=text,
                text_extracted=bool(text.strip()),
            )
        )
    return slides


# --------------------------------------------------------------------------- #
# PPTX
# --------------------------------------------------------------------------- #


def _shape_texts(shape) -> list[str]:  # noqa: ANN001 - python-pptx shapes are untyped
    """Pull every readable string out of a single shape, recursing into groups."""
    texts: list[str] = []

    # Grouped shapes hold their own children.
    if getattr(shape, "shape_type", None) is not None and hasattr(shape, "shapes"):
        for child in shape.shapes:
            texts.extend(_shape_texts(child))
        return texts

    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            line = "".join(run.text for run in paragraph.runs)
            if line.strip():
                texts.append(line)

    if getattr(shape, "has_table", False):
        table = shape.table
        for row_idx, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                prefix = "Table header: " if row_idx == 0 else ""
                texts.append(prefix + " | ".join(cells))

    if getattr(shape, "has_chart", False):
        chart = shape.chart
        try:
            if chart.has_title and chart.chart_title.has_text_frame:
                texts.append(f"Chart: {chart.chart_title.text_frame.text}")
        except Exception:  # noqa: BLE001 - chart titles are frequently malformed
            pass
        try:
            for series in chart.plots[0].categories:
                if series:
                    texts.append(f"Chart category: {series}")
        except Exception:  # noqa: BLE001
            pass

    return texts


def _extract_pptx(path: Path) -> list[SlideContent]:
    from pptx import Presentation

    try:
        presentation = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001 - python-pptx raises package/XML errors
        raise ExtractionError(f"Could not open PPTX '{path.name}': {exc}") from exc

    slides: list[SlideContent] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            try:
                texts.extend(_shape_texts(shape))
            except Exception:  # noqa: BLE001 - skip shapes that fail to render text
                continue

        title = None
        try:
            if slide.shapes.title is not None and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
        except Exception:  # noqa: BLE001
            title = None

        text = clean_text("\n".join(texts))
        slides.append(
            SlideContent(
                slide_number=index,
                title=title or guess_title(text),
                raw_text=text,
                text_extracted=bool(text.strip()),
            )
        )
    return slides


__all__ = ["extract_deck", "ExtractionError", "FileError"]
