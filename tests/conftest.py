"""Shared fixtures: synthetic decks and a canned analysis."""

from __future__ import annotations

import json
from pathlib import Path

import pytest

from pitch2onepager.models import CustomerJourneyAnalysis

# (title, body, footer) — the footer is drawn on its own line so the fixtures
# exercise the real boilerplate-stripping rules.
SLIDE_TEXT = [
    (
        "Meridian Health",
        "Meridian Health closes the gap between hospital discharge and home recovery.",
        "Confidential",
    ),
    (
        "The Problem",
        "Care coordinators discover a readmission risk only after the patient is back "
        "in the ER. 1 in 5 Medicare patients is readmitted within 30 days, costing "
        "US hospitals $26B annually. Coordinators wait 3 to 6 weeks before escalating "
        "to a vendor search. Contact hello@meridianhealth.example for details.",
        "Page 2",
    ),
    (
        "What Exists Today",
        "Hospitals evaluate 5 to 10 vendors over 4 to 9 months. They find them through "
        "Google search, peer referrals from other health systems, KLAS reports, and "
        "HIMSS conference floors. Legacy discharge modules inside the EHR do not "
        "surface risk. Point solutions do not write back to the chart. Staffing "
        "agencies cost more than the readmission penalty they prevent.",
        "© 2026 Meridian Health",
    ),
    (
        "Workarounds",
        "Coordinators keep parallel spreadsheets, run manual call-down lists every "
        "Friday, and text nurse managers after hours. One director described it as "
        "'flying the plane with a paper map'.",
        "4",
    ),
    (
        "Why Now",
        "CMS penalty tiers expand in 2026 and FHIR write-back became mandatory for "
        "certified EHRs. TAM is $4.2B across 6,100 US acute-care hospitals.",
        "Page 5",
    ),
]


def _analysis_dict() -> dict:
    return {
        "company_name": "Meridian Health",
        "tagline": "Closing the gap between hospital discharge and home recovery.",
        "target_customer": "Care coordination directors at US acute-care hospitals",
        "problem_awareness": {
            "trigger_moment": (
                "A coordinator learns a discharged patient is back in the ER — the "
                "first signal the handoff failed."
            ),
            "pain_description": (
                "1 in 5 Medicare patients is readmitted within 30 days, costing US "
                "hospitals $26B annually and triggering CMS penalties the coordinator "
                "is personally measured on."
            ),
            "time_before_seeking_solution": "3–6 weeks",
            "urgency_level": "high",
        },
        "discovery_odyssey": {
            "channels_navigated": [
                "Google search",
                "Peer referrals",
                "KLAS reports",
                "HIMSS conference",
            ],
            "vendor_count_estimate": "5–10 vendors evaluated",
            "avg_time_to_find_fit": "4–9 months",
            "friction_points": [
                "No neutral comparison source for post-acute vendors",
                "EHR integration claims are unverifiable before contract",
                "Committee sign-off spans clinical, IT, and finance",
            ],
        },
        "solution_landscape": {
            "existing_solutions": [
                "Legacy EHR discharge modules",
                "Point-solution readmission apps",
                "Staffing agencies",
            ],
            "key_shortfalls": [
                "Do not surface risk before discharge",
                "No write-back to the chart",
                "Cost exceeds the penalty they prevent",
            ],
            "workarounds_customers_use": [
                "Parallel spreadsheets outside the EHR",
                "Manual Friday call-down lists",
                "After-hours texts to nurse managers",
            ],
        },
        "day_in_the_life": {
            "operational_burden": (
                "Coordinators rebuild the same at-risk list by hand every week across "
                "three disconnected systems."
            ),
            "emotional_burden": (
                "Directors carry personal accountability for penalties they lack the "
                "data to prevent."
            ),
            "financial_burden": (
                "A single hospital absorbs $1.2M–$4M in annual penalties plus the "
                "unreimbursed cost of each readmission."
            ),
            "organizational_spillover": (
                "Nursing, IT, and finance are each pulled into escalations that begin "
                "as a coordination gap."
            ),
        },
        "gaps_that_matter": {
            "unmet_needs": [
                "Risk signal before discharge, not after readmission",
                "Write-back into the chart of record",
                "Pricing below the penalty being avoided",
            ],
            "customer_quotes": ["flying the plane with a paper map"],
        },
        "investment_thesis": {
            "market_opportunity_statement": (
                "Post-acute coordination is a $4.2B market where the buyer is measured "
                "on an outcome no incumbent gives them the data to control. CMS penalty "
                "expansion in 2026 converts a soft operational frustration into a hard "
                "budget line, and mandatory FHIR write-back removes the integration "
                "barrier that protected legacy EHR modules."
            ),
            "comparable_industries": [
                "Enterprise HR tech before Workday",
                "Healthcare data interoperability pre-2015",
            ],
            "why_now": "CMS penalty tiers expand in 2026 and FHIR write-back is mandatory",
            "estimated_market_size": "$4.2B TAM across 6,100 US acute-care hospitals",
        },
    }


@pytest.fixture
def analysis_json() -> str:
    return json.dumps(_analysis_dict())


@pytest.fixture
def analysis() -> CustomerJourneyAnalysis:
    return CustomerJourneyAnalysis.model_validate(_analysis_dict())


@pytest.fixture
def sample_pdf(tmp_path: Path) -> Path:
    """A synthetic, text-based pitch deck PDF built with reportlab."""
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    path = tmp_path / "sample_deck.pdf"
    c = canvas.Canvas(str(path), pagesize=letter)
    for title, body, footer in SLIDE_TEXT:
        c.setFont("Helvetica-Bold", 20)
        c.drawString(60, 700, title)
        c.setFont("Helvetica", 11)
        y = 660
        words, line = body.split(), ""
        for word in words:
            if len(line) + len(word) + 1 > 78:
                c.drawString(60, y, line)
                y -= 16
                line = word
            else:
                line = f"{line} {word}".strip()
        if line:
            c.drawString(60, y, line)
        c.setFont("Helvetica", 8)
        c.drawString(60, 60, footer)
        c.showPage()
    c.save()
    return path


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """A synthetic pitch deck PPTX built with python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    path = tmp_path / "sample_deck.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[6]

    for title, body, footer in SLIDE_TEXT:
        slide = prs.slides.add_slide(blank)
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].runs[0].font.size = Pt(28)

        body_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.6), Inches(9), Inches(4)
        )
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        body_frame.text = body

        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.8), Inches(9), Inches(0.4)
        )
        footer_box.text_frame.text = footer

    # A table on the last slide exercises the table-header extraction path.
    last = prs.slides[-1]
    table = last.shapes.add_table(2, 2, Inches(0.5), Inches(5.5), Inches(6), Inches(1)).table
    table.cell(0, 0).text = "Segment"
    table.cell(0, 1).text = "TAM"
    table.cell(1, 0).text = "Acute care"
    table.cell(1, 1).text = "$4.2B"

    prs.save(str(path))
    return path


@pytest.fixture
def image_only_pdf(tmp_path: Path) -> Path:
    """A PDF whose pages carry no extractable text."""
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    path = tmp_path / "image_only.pdf"
    c = canvas.Canvas(str(path), pagesize=letter)
    for _ in range(3):
        c.rect(100, 400, 300, 200, stroke=1, fill=0)
        c.showPage()
    c.save()
    return path


class FakeTextBlock:
    """Mimics an anthropic TextBlock."""

    def __init__(self, text: str) -> None:
        self.type = "text"
        self.text = text


class FakeMessage:
    def __init__(self, text: str) -> None:
        self.content = [FakeTextBlock(text)]


class FakeMessages:
    def __init__(self, responses: list[str]) -> None:
        self._responses = list(responses)
        self.calls: list[dict] = []

    def create(self, **kwargs):  # noqa: ANN003, ANN201
        self.calls.append(kwargs)
        if not self._responses:
            raise AssertionError("FakeMessages.create called more times than expected")
        return FakeMessage(self._responses.pop(0))


class FakeAnthropic:
    """Stand-in for ``anthropic.Anthropic`` that replays canned responses."""

    def __init__(self, responses: list[str]) -> None:
        self.messages = FakeMessages(responses)


@pytest.fixture
def fake_client_factory():
    def _make(responses: list[str]) -> FakeAnthropic:
        return FakeAnthropic(responses)

    return _make
