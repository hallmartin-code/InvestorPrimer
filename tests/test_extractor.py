from __future__ import annotations

from pathlib import Path

import pytest

from pitch2onepager.extractor import extract_deck
from pitch2onepager.utils import (
    ExtractionError,
    FileError,
    clean_text,
    detect_file_type,
    guess_title,
    slugify,
    truncate,
)


class TestFileTypeDetection:
    def test_detects_pdf_and_pptx(self, tmp_path: Path) -> None:
        assert detect_file_type(tmp_path / "a.pdf") == "pdf"
        assert detect_file_type(tmp_path / "a.PPTX") == "pptx"

    def test_rejects_unsupported_extension(self, tmp_path: Path) -> None:
        with pytest.raises(FileError, match="Supported formats"):
            detect_file_type(tmp_path / "deck.key")

    def test_missing_file_raises(self, tmp_path: Path) -> None:
        with pytest.raises(FileError, match="File not found"):
            extract_deck(str(tmp_path / "nope.pdf"))

    def test_directory_raises(self, tmp_path: Path) -> None:
        d = tmp_path / "adir.pdf"
        d.mkdir()
        with pytest.raises(FileError, match="Not a file"):
            extract_deck(str(d))


class TestPdfExtraction:
    def test_extracts_every_page(self, sample_pdf: Path) -> None:
        deck = extract_deck(str(sample_pdf))
        assert deck.file_type == "pdf"
        assert deck.slide_count == 5
        assert deck.extracted_slide_count == 5
        assert [s.slide_number for s in deck.slides] == [1, 2, 3, 4, 5]

    def test_full_text_contains_key_claims(self, sample_pdf: Path) -> None:
        deck = extract_deck(str(sample_pdf))
        assert "readmitted within 30 days" in deck.full_text
        assert "Meridian Health" in deck.full_text

    def test_titles_are_guessed(self, sample_pdf: Path) -> None:
        deck = extract_deck(str(sample_pdf))
        assert deck.slides[1].title == "The Problem"

    def test_boilerplate_is_stripped(self, sample_pdf: Path) -> None:
        deck = extract_deck(str(sample_pdf))
        assert "hello@meridianhealth.example" not in deck.full_text
        assert "Confidential" not in deck.full_text
        assert "Page 5" not in deck.full_text
        assert "© 2026" not in deck.full_text
        # Surrounding prose survives the inline email removal.
        assert "Contact for details." in deck.full_text

    def test_image_only_pages_marked_not_extracted(self, image_only_pdf: Path) -> None:
        deck = extract_deck(str(image_only_pdf))
        assert deck.slide_count == 3
        assert deck.extracted_slide_count == 0
        assert all(s.text_extracted is False for s in deck.slides)

    def test_corrupt_pdf_raises_extraction_error(self, tmp_path: Path) -> None:
        bad = tmp_path / "broken.pdf"
        bad.write_bytes(b"this is definitely not a pdf")
        with pytest.raises(ExtractionError, match="Could not open PDF"):
            extract_deck(str(bad))


class TestPptxExtraction:
    def test_extracts_every_slide(self, sample_pptx: Path) -> None:
        deck = extract_deck(str(sample_pptx))
        assert deck.file_type == "pptx"
        assert deck.slide_count == 5
        assert deck.extracted_slide_count == 5

    def test_captures_body_and_table_content(self, sample_pptx: Path) -> None:
        deck = extract_deck(str(sample_pptx))
        assert "5 to 10 vendors" in deck.full_text
        assert "Table header: Segment | TAM" in deck.full_text
        assert "Acute care | $4.2B" in deck.full_text

    def test_extracts_grouped_shapes_and_chart_titles(self, tmp_path: Path) -> None:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches

        path = tmp_path / "rich.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        group = slide.shapes.add_group_shape()
        for i, label in enumerate(("Grouped headline", "Grouped detail line")):
            box = group.shapes.add_textbox(
                Inches(0.5), Inches(0.5 + i), Inches(4), Inches(0.8)
            )
            box.text_frame.text = label

        chart_data = CategoryChartData()
        chart_data.categories = ["Acute care", "Post acute"]
        chart_data.add_series("TAM", (4.2, 1.8))
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.5), Inches(3.5), Inches(6), Inches(3),
            chart_data,
        )
        chart = graphic_frame.chart
        chart.has_title = True
        chart.chart_title.text_frame.text = "Market by segment"

        prs.save(str(path))

        deck = extract_deck(str(path))
        text = deck.full_text
        assert "Grouped headline" in text
        assert "Grouped detail line" in text
        assert "Chart: Market by segment" in text
        assert "Chart category: Acute care" in text

    def test_corrupt_pptx_raises_extraction_error(self, tmp_path: Path) -> None:
        bad = tmp_path / "broken.pptx"
        bad.write_bytes(b"not a zip archive")
        with pytest.raises(ExtractionError, match="Could not open PPTX"):
            extract_deck(str(bad))


class TestTextHelpers:
    @pytest.mark.parametrize(
        "line",
        ["  ", "12", "Page 7", "3 / 20", "Confidential", "© 2026 Acme Inc", "www@x.co"],
    )
    def test_boilerplate_lines_dropped(self, line: str) -> None:
        assert clean_text(f"Keep this line\n{line}") == "Keep this line"

    def test_company_hint_line_dropped(self) -> None:
        assert clean_text("Acme\nReal content", company_hint="acme") == "Real content"

    def test_whitespace_collapsed(self) -> None:
        assert clean_text("a    b\t\tc") == "a b c"

    def test_empty_input(self) -> None:
        assert clean_text("") == ""

    def test_guess_title_skips_blank_lines(self) -> None:
        assert guess_title("\n\n  Real Title  \nbody") == "Real Title"
        assert guess_title("   ") is None

    def test_truncate(self) -> None:
        assert truncate("short", 20) == "short"
        assert truncate("abcdefghij", 5).endswith("…")
        assert len(truncate("abcdefghij", 5)) == 5

    def test_slugify(self) -> None:
        assert slugify("Meridian Health, Inc.") == "Meridian_Health_Inc"
        assert slugify("///") == "company"
